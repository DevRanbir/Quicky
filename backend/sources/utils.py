import PyPDF2
import docx
from pptx import Presentation
from youtube_transcript_api import YouTubeTranscriptApi
from youtube_transcript_api import (
    NoTranscriptFound,
    TranscriptsDisabled,
    NoTranscriptAvailable,
)

import re

def extract_text_from_pdf(file_path):
    page_texts = []
    page_count = 0
    try:
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            page_count = len(reader.pages)
            for page_num in range(page_count):
                page_text = reader.pages[page_num].extract_text()
                if page_text: # Ensure there's text to add
                    page_texts.append(page_text)
                else:
                    page_texts.append("") # Add empty string for blank pages to maintain page count integrity
    except Exception as e:
        print(f"Error extracting PDF: {e}")
        return None, 0 # Return None for texts and 0 for count on error
    return page_texts, page_count

def extract_text_from_docx(file_path):
    page_texts = []
    current_page_paragraphs = []
    PARAGRAPHS_PER_PAGE = 10  # Define how many paragraphs constitute a "page"
    paragraph_count_on_current_page = 0

    try:
        doc = docx.Document(file_path)
        # First check if document has any text content
        has_content = False
        for para in doc.paragraphs:
            if para.text.strip():  # Check for non-empty paragraphs
                has_content = True
                break
        
        if not has_content:
            print(f"Warning: DOCX file {file_path} appears to be empty")
            return None, 0

        for para in doc.paragraphs:
            text = para.text.strip()
            if text:  # Only add non-empty paragraphs
                current_page_paragraphs.append(text)
                paragraph_count_on_current_page += 1
                if paragraph_count_on_current_page >= PARAGRAPHS_PER_PAGE:
                    page_text = "\n".join(current_page_paragraphs)
                    if page_text.strip():  # Only add non-empty pages
                        page_texts.append(page_text)
                    current_page_paragraphs = []
                    paragraph_count_on_current_page = 0
        
        # Add any remaining paragraphs as the last page
        if current_page_paragraphs:
            page_text = "\n".join(current_page_paragraphs)
            if page_text.strip():  # Only add non-empty pages
                page_texts.append(page_text)
        
        # If we haven't created any pages but have content, create a single page
        if not page_texts:
            all_text = "\n".join([p.text.strip() for p in doc.paragraphs if p.text.strip()])
            if all_text:
                page_texts.append(all_text)

        return page_texts, len(page_texts) if page_texts else 0
    except Exception as e:
        print(f"Error extracting DOCX: {e}")
        return None, 0

def extract_text_from_pptx(file_path):
    slide_texts = []
    try:
        prs = Presentation(file_path)
        
        # First check if presentation has any text content
        has_content = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    has_content = True
                    break
            if has_content:
                break
        
        if not has_content:
            print(f"Warning: PPTX file {file_path} appears to be empty")
            return None, 0

        for slide in prs.slides:
            slide_text_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:  # Only add non-empty text
                        slide_text_content.append(text)
            
            if slide_text_content:  # Add slide only if it has text
                slide_text = "\n".join(slide_text_content)
                if slide_text.strip():  # Double check the joined text is not empty
                    slide_texts.append(slide_text)
                else:
                    slide_texts.append("")  # Maintain slide count integrity
            else:
                slide_texts.append("")  # Maintain slide count integrity
        
        # If we haven't found any text content but have slides, return None
        if not any(text.strip() for text in slide_texts):
            print(f"Warning: No text content found in PPTX file {file_path}")
            return None, 0
            
        return slide_texts, len(prs.slides)
    except Exception as e:
        print(f"Error extracting PPTX: {e}")
        return None, 0

def extract_text_from_txt(file_path):
    text = ""
    encodings = ['utf-8', 'ascii', 'iso-8859-1', 'cp1252', 'utf-16']
    
    for encoding in encodings:
        try:
            with open(file_path, 'r', encoding=encoding) as f:
                text = f.read().strip()
                if not text:  # Check if file is empty
                    print(f"Warning: TXT file {file_path} appears to be empty")
                    return None
                return text
        except UnicodeDecodeError:
            continue  # Try next encoding
        except Exception as e:
            print(f"Error extracting TXT with {encoding} encoding: {e}")
            continue
    
    print(f"Error: Could not read TXT file {file_path} with any supported encoding")
    return None

def extract_youtube_id(youtube_url):
    regex = r"(?:https?:\/\/)?(?:www\.)?(?:youtube\.com\/(?:[^\/\n\s]+\/\S+\/|[^\/\n\s]+\/\S+\/|(?:v|e(?:mbed)?)\/|\S*?[?&]v=)|youtu\.be\/)([a-zA-Z0-9_-]{11})"
    match = re.search(regex, youtube_url)
    return match.group(1) if match else None

def extract_youtube_transcript(video_id):
    """Extracts YouTube video transcript, preferring English or translating if necessary."""
    print(f"Attempting to extract transcript for video ID: {video_id}")
    
    # Initialize with default values
    title = ""
    channel = ""
    duration = ""
    thumbnail = ""
    transcript_list = [] # Initialize transcript_list here
    
    try:
        # First try to get English transcript
        try:
            transcript_list = YouTubeTranscriptApi.get_transcript(video_id, languages=['en'])
        except (NoTranscriptFound, TranscriptsDisabled, NoTranscriptAvailable) as e:
            print(f"English transcript not found or disabled for {video_id}: {e}")
            # If English not available, try to get any available transcript and translate to English
            try:
                all_transcripts = YouTubeTranscriptApi.list_transcripts(video_id)
                available_transcript = None
                
                # Try to find a generated transcript
                try:
                    available_transcript = all_transcripts.find_generated_transcript(['en', 'hi'])
                except NoTranscriptFound:
                    print(f"No generated transcript found for {video_id}.")
                
                # If no generated transcript, try manually created ones
                if not available_transcript:
                    try:
                        available_transcript = all_transcripts.find_manually_created_transcript(['en'])
                    except NoTranscriptFound:
                        print(f"No manually created transcript found for {video_id}.")
                
                # If we found a transcript, fetch and translate it if needed
                if available_transcript:
                    if available_transcript.language_code != 'en':
                        transcript_list = available_transcript.translate('en').fetch()
                    else:
                        transcript_list = available_transcript.fetch()
                else:
                    # Last resort: try to get any transcript and translate
                    try:
                        # Get the first available transcript, regardless of language
                        first_transcript = next(iter(all_transcripts), None)
                        if first_transcript:
                            transcript_list = first_transcript.translate('en').fetch()
                        else:
                            print(f"No transcripts available at all for {video_id}.")
                            transcript_list = [] # Ensure it's empty if nothing found
                    except Exception as e3:
                        print(f"Failed to get any transcript for {video_id}: {e3}")
                        transcript_list = [] # Ensure it's empty on error
            except Exception as outer_e:
                print(f"Failed to list or process transcripts for {video_id}: {outer_e}")
                transcript_list = [] # Ensure it's empty on error
        
        # Process the transcript
        if not transcript_list:
            print(f"No transcript items found for video_id: {video_id}")
            return None

        text = " "
        for item in transcript_list:
            text += item['text'] + " "
        
        return text.strip()
    except Exception as e:
        print(f"Error extracting YouTube transcript for {video_id}: {e}")
        return None, None