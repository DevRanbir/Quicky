from rest_framework import viewsets, status
from rest_framework.response import Response
from rest_framework.decorators import action
from .models import Source
from .serializers import FileUploadSerializer, YouTubeLinkSerializer, SourceSerializer
from .utils import (
    extract_text_from_pdf, extract_text_from_docx, 
    extract_text_from_pptx, extract_text_from_txt,
    extract_youtube_transcript, extract_youtube_id
)
from django.core.files.storage import default_storage
import os
from django.conf import settings

import json
from django.http import JsonResponse
from django.views.decorators.http import require_http_methods
from django.views.decorators.csrf import csrf_exempt
from rest_framework.decorators import api_view
import PyPDF2
import docx
from pptx import Presentation
import yt_dlp
import requests
from urllib.parse import urlparse, parse_qs
import re
from youtube_transcript_api import YouTubeTranscriptApi
from youtube_transcript_api.formatters import TextFormatter

# Import the actual function - remove the try/except wrapper
from questions.utils import generate_questions_from_text_content



class SourceViewSet(viewsets.ModelViewSet):
    queryset = Source.objects.all().order_by('-uploaded_at')
    serializer_class = SourceSerializer

    @action(detail=False, methods=['post'], serializer_class=FileUploadSerializer)
    def upload_file(self, request):
        serializer = self.get_serializer(data=request.data)
        if serializer.is_valid():
            uploaded_file = serializer.validated_data['file']

            # Check for duplicate file name
            if Source.objects.filter(file=uploaded_file.name).exists():
                return Response({"error": f"File with name '{uploaded_file.name}' already exists."}, status=status.HTTP_409_CONFLICT)

            file_name = default_storage.save(uploaded_file.name, uploaded_file)
            file_path = default_storage.path(file_name)
            
            text_content = None
            page_count = None
            source_type_enum = None

            ext = file_name.split('.')[-1].lower()
            if ext == 'pdf':
                text_content, page_count = extract_text_from_pdf(file_path)
                source_type_enum = 'PDF'
            elif ext == 'docx':
                text_content, page_count = extract_text_from_docx(file_path)
                source_type_enum = 'DOCX'
            elif ext == 'pptx':
                text_content, page_count = extract_text_from_pptx(file_path)
                source_type_enum = 'PPTX'
            elif ext == 'txt':
                text_content = extract_text_from_txt(file_path)
                source_type_enum = 'TXT'

            if text_content is not None:
                # For non-PDFs, wrap the text_content string in a list to match JSONField structure
                # For PDFs, DOCX, PPTX, text_content is already a list of strings (pages/slides)
                # For TXT, wrap the text_content string in a list to match JSONField structure
                if source_type_enum == 'TXT' and isinstance(text_content, str):
                    processed_text_content = [text_content]
                else:
                    processed_text_content = text_content

                source = Source.objects.create(
                    source_type=source_type_enum,
                    file=file_name, 
                    text_content=processed_text_content,
                    page_count=page_count
                )
                return Response(SourceSerializer(source, context={'request': request}).data, status=status.HTTP_201_CREATED)
            else:
                # If text extraction failed, delete the uploaded file
                if default_storage.exists(file_name):
                    default_storage.delete(file_name)
                return Response({"error": "Failed to extract text from file."}, status=status.HTTP_400_BAD_REQUEST)
        return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)

    @action(detail=False, methods=['post'], serializer_class=YouTubeLinkSerializer)
    def process_youtube_link(self, request):
        serializer = self.get_serializer(data=request.data)
        if serializer.is_valid():
            youtube_link = serializer.validated_data['youtube_link']
            
            # Check for duplicate YouTube link
            if Source.objects.filter(youtube_link=youtube_link).exists():
                return Response({"error": f"YouTube link '{youtube_link}' has already been processed."}, status=status.HTTP_409_CONFLICT)

            video_id = extract_youtube_id(youtube_link)
            if not video_id:
                return Response({"error": "Invalid YouTube URL or could not extract video ID."}, status=status.HTTP_400_BAD_REQUEST)
            
            text_content, video_duration = extract_youtube_transcript(video_id)
            
            if text_content:
                # Wrap the YouTube transcript in a list to match JSONField structure
                processed_text_content = [text_content] 
                source = Source.objects.create(
                    source_type='YOUTUBE',
                    youtube_link=youtube_link,
                    text_content=processed_text_content,
                    video_duration=video_duration
                )
                return Response(SourceSerializer(source, context={'request': request}).data, status=status.HTTP_201_CREATED)
            else:
                return Response({"error": "Failed to extract transcript from YouTube video."}, status=status.HTTP_400_BAD_REQUEST)
        return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)

    @action(detail=False, methods=['get'])
    def files(self, request):
        sources = Source.objects.all().order_by('-uploaded_at')
        serializer = SourceSerializer(sources, many=True, context={'request': request})
        return Response(serializer.data)

    @action(detail=True, methods=['delete'])
    def delete_file(self, request, pk=None):
        try:
            source = self.get_object()
            if source.file and default_storage.exists(source.file.name):
                default_storage.delete(source.file.name)
            source.delete()
            return Response(status=status.HTTP_204_NO_CONTENT)
        except Source.DoesNotExist:
            return Response({"error": "File not found."}, status=status.HTTP_404_NOT_FOUND)
        except Exception as e:
            return Response({"error": str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

    @action(detail=True, methods=['post'])
    def generate_questions(self, request, pk=None):
        try:
            source = self.get_object()
        except Source.DoesNotExist:
            return Response({"error": "Source not found."}, status=status.HTTP_404_NOT_FOUND)

        pages_to_generate_str = request.data.get('pages_to_generate') # e.g., "1-5,7,10-12" or empty for all/non-PDF
        questions_per_page_str = request.data.get('questions_per_page', '5') # Default to 5 questions per page
        total_question_limit_str = request.data.get('total_question_limit') # Optional overall limit

        try:
            questions_per_page = int(questions_per_page_str)
            if questions_per_page <= 0:
                return Response({"error": "questions_per_page must be a positive integer."}, status=status.HTTP_400_BAD_REQUEST)
        except ValueError:
            return Response({"error": "Invalid questions_per_page. Must be an integer."}, status=status.HTTP_400_BAD_REQUEST)
        
        total_question_limit = None
        if total_question_limit_str:
            try:
                total_question_limit = int(total_question_limit_str)
                if total_question_limit <= 0:
                    return Response({"error": "total_question_limit must be a positive integer if provided."}, status=status.HTTP_400_BAD_REQUEST)
            except ValueError:
                return Response({"error": "Invalid total_question_limit. Must be an integer."}, status=status.HTTP_400_BAD_REQUEST)

        # Store generation parameters in source_metadata
        source.source_metadata = {
            'pages_to_generate': pages_to_generate_str,
            'questions_per_page': questions_per_page,
            'total_question_limit': total_question_limit
        }
        source.save() # Save metadata

        # Ensure source.text_content is available and is a list (as expected by the util)
        if not source.text_content or not isinstance(source.text_content, list):
             # This might happen if a very old source record didn't get its text_content as a list
             # Or if text extraction failed previously but somehow this endpoint is hit.
            return Response({"error": "Source content is not available or not in the expected format (list of page texts)."}, status=status.HTTP_400_BAD_REQUEST)

        try:
            generated_questions_data = generate_questions_from_text_content(
                source_text_content=source.text_content, # This is now a list of texts per page for PDF
                questions_per_page=questions_per_page,
                pages_to_generate_str=pages_to_generate_str,
                total_question_limit=total_question_limit,
                source_id=source.id  # Added missing source_id parameter
            )

            if generated_questions_data:
                return Response(generated_questions_data, status=status.HTTP_200_OK)
            else:
                # The utility function now prints more specific warnings, so a generic error here is okay.
                return Response({"error": "Failed to generate questions. This could be due to empty content on specified pages, invalid page ranges, or an issue with the content processing."}, status=status.HTTP_400_BAD_REQUEST)
        
        except Exception as e:
            print(f"Error in generate_questions endpoint: {str(e)}")
            return Response({"error": f"Internal server error during question generation: {str(e)}"}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
            
    @action(detail=True, methods=['get'])
    def file(self, request, pk=None):
        try:
            source = self.get_object()
            if not source.file:
                return Response({"error": "No file associated with this source."}, status=status.HTTP_404_NOT_FOUND)
            
            serializer = SourceSerializer(source, context={'request': request})
            return Response({"file_url": serializer.data['file_url']})
        except Source.DoesNotExist:
            return Response({"error": "Source not found."}, status=status.HTTP_404_NOT_FOUND)
        except Exception as e:
            return Response({"error": str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

@api_view(['GET'])
def source_preview(request, source_id):
    """
    Generate preview content for a source file
    """
    try:
        # Get the source object
        source = Source.objects.get(id=source_id)
        
        # Get optional page limit from query params (default 10 for better performance)
        page_limit = int(request.GET.get('page_limit', 100))
        
        preview_data = {}
        
        if source.source_type == 'PDF':
            preview_data = generate_pdf_preview(source, page_limit)
        elif source.source_type == 'DOCX':
            preview_data = generate_docx_preview(source)
        elif source.source_type == 'PPTX':
            preview_data = generate_pptx_preview(source)
        elif source.source_type == 'TXT':
            preview_data = generate_txt_preview(source)
        elif source.source_type == 'YOUTUBE':
            preview_data = generate_youtube_preview(source)
        else:
            return Response(
                {'error': f'Preview not supported for {source.source_type}'}, 
                status=status.HTTP_400_BAD_REQUEST
            )
            
        return Response(preview_data, status=status.HTTP_200_OK)
        
    except Source.DoesNotExist:
        return Response(
            {'error': 'Source not found'}, 
            status=status.HTTP_404_NOT_FOUND
        )
    except Exception as e:
        print(f"Preview error for source {source_id}: {str(e)}")
        return Response(
            {'error': f'Failed to generate preview: {str(e)}'}, 
            status=status.HTTP_500_INTERNAL_SERVER_ERROR
        )

def generate_pdf_preview(source, page_limit=100):
    """Generate preview for PDF files with improved performance"""
    try:
        # Use text_content from database if available (much faster)
        if source.text_content and isinstance(source.text_content, list):
            pages = []
            total_pages = len(source.text_content)
            
            # Show up to page_limit pages
            for i, page_text in enumerate(source.text_content[:page_limit]):
                if page_text and page_text.strip():
                    # Clean up the text
                    cleaned_text = re.sub(r'\s+', ' ', page_text).strip()
                    pages.append({
                        'page_number': i + 1,
                        'content': cleaned_text[:1000] if len(cleaned_text) > 1000 else cleaned_text,
                        'word_count': len(cleaned_text.split())
                    })
            
            return {
                'pages': pages,
                'total_pages': total_pages,
                'preview_pages': min(page_limit, total_pages),
                'source': 'database'
            }
        
        # Fallback to file reading if text_content is not available
        if not source.file:
            raise Exception("No file path available")
            
        file_path = source.file.path
        pages = []
        
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            total_pages = len(pdf_reader.pages)
            
            # Limit preview to specified page_limit or all pages if less
            max_pages = min(total_pages, page_limit)
            
            for i in range(max_pages):
                try:
                    page = pdf_reader.pages[i]
                    text = page.extract_text()
                    
                    if text and text.strip():
                        # Clean up the text
                        text = re.sub(r'\s+', ' ', text).strip()
                        
                        pages.append({
                            'page_number': i + 1,
                            'content': text[:1000] if len(text) > 1000 else text,
                            'word_count': len(text.split())
                        })
                except Exception as page_error:
                    print(f"Error reading page {i+1}: {str(page_error)}")
                    pages.append({
                        'page_number': i + 1,
                        'content': f"Error reading page {i+1}",
                        'word_count': 0
                    })
        
        return {
            'pages': pages,
            'total_pages': total_pages,
            'preview_pages': min(page_limit, total_pages),
            'source': 'file'
        }
        
    except Exception as e:
        print(f"PDF preview error: {str(e)}")
        raise Exception(f'Error reading PDF: {str(e)}')

def generate_docx_preview(source):
    """Generate preview for DOCX files with improved performance"""
    try:
        # Use text_content from database if available
        if source.text_content and isinstance(source.text_content, list) and len(source.text_content) > 0:
            full_text = source.text_content[0]  # DOCX is stored as single text block
            
            # Split into paragraphs for analysis
            paragraphs = [p.strip() for p in full_text.split('\n') if p.strip()]
            
            return {
                'text': full_text[:200000] if len(full_text) > 200000 else full_text,
                'paragraph_count': len(paragraphs),
                'word_count': len(full_text.split()),
                'character_count': len(full_text),
                'source': 'database'
            }
        
        # Fallback to file reading
        if not source.file:
            raise Exception("No file path available")
            
        file_path = source.file.path
        doc = docx.Document(file_path)
        
        # Extract text from all paragraphs
        text_content = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text.strip())
        
        full_text = '\n'.join(text_content)
        
        return {
            'text_content': full_text[:200000] if len(full_text) > 200000 else full_text,
            'paragraph_count': len(text_content),
            'word_count': len(full_text.split()),
            'character_count': len(full_text),
            'source': 'file'
        }
        
    except Exception as e:
        print(f"DOCX preview error: {str(e)}")
        raise Exception(f'Error reading DOCX: {str(e)}')

def generate_pptx_preview(source):
    """Generate preview for PPTX files with improved performance"""
    try:
        # Use text_content from database if available
        if source.text_content and isinstance(source.text_content, list) and len(source.text_content) > 0:
            full_text = source.text_content[0]  # PPTX is stored as single text block
            
            # Try to split into slides if possible
            slides_text = full_text.split('\n\n')  # Rough slide separation
            
            return {
                'text': full_text[:200000] if len(full_text) > 200000 else full_text,
                'word_count': len(full_text.split()),
                'character_count': len(full_text),
                'estimated_slides': len([s for s in slides_text if s.strip()]),
                'source': 'database'
            }
        
        # Fallback to file reading
        if not source.file:
            raise Exception("No file path available")
            
        file_path = source.file.path
        prs = Presentation(file_path)
        
        slides_content = []
        
        for i, slide in enumerate(prs.slides):
            slide_text = []
            
            # Extract text from all shapes in slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                slides_content.append({
                    'slide_number': i + 1,
                    'content': '\n'.join(slide_text)[:50000]  # Limit per slide
                })
        
        # Create a combined text preview
        all_text = []
        for slide in slides_content:
            all_text.append(f"Slide {slide['slide_number']}:\n{slide['content']}")
        
        combined_text = '\n\n'.join(all_text)
        
        return {
            'text_content': combined_text[:200000] if len(combined_text) > 200000 else combined_text,
            'slides': slides_content[:5],  # Show first 5 slides in detail
            'total_slides': len(prs.slides),
            'word_count': len(combined_text.split()),
            'source': 'file'
        }
        
    except Exception as e:
        print(f"PPTX preview error: {str(e)}")
        raise Exception(f'Error reading PPTX: {str(e)}')

def generate_txt_preview(source):
    """Generate preview for TXT files with improved performance"""
    try:
        # Use text_content from database if available
        if source.text_content and isinstance(source.text_content, list) and len(source.text_content) > 0:
            content = source.text_content[0]  # TXT is stored as single text block
            
            lines = content.split('\n')
            words = content.split()
            
            return {
                'text_content': content[:200000] if len(content) > 200000 else content,
                'line_count': len(lines),
                'word_count': len(words),
                'character_count': len(content),
                'source': 'database'
            }
        
        # Fallback to file reading
        if not source.file:
            raise Exception("No file path available")
            
        file_path = source.file.path
        
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
        except UnicodeDecodeError:
            # Try with different encoding
            with open(file_path, 'r', encoding='latin-1') as file:
                content = file.read()
        
        lines = content.split('\n')
        words = content.split()
        
        return {
            'text_content': content[:200000] if len(content) > 200000 else content,
            'line_count': len(lines),
            'word_count': len(words),
            'character_count': len(content),
            'source': 'file'
        }
        
    except Exception as e:
        print(f"TXT preview error: {str(e)}")
        raise Exception(f'Error reading TXT: {str(e)}')

def generate_youtube_preview(source):
    """Generate preview for YouTube videos with improved transcript handling"""
    try:
        video_url = source.youtube_link
        
        # Extract video ID from URL
        video_id = extract_youtube_video_id(video_url)
        if not video_id:
            raise Exception('Invalid YouTube URL')
        
        # Use text_content from database if available (much faster)
        transcript_text = ""
        if source.text_content and isinstance(source.text_content, list) and len(source.text_content) > 0:
            transcript_text = source.text_content[0][:150000]  # Limit for preview
            
        # Get basic video information using yt-dlp (lightweight extraction)
        try:
            ydl_opts = {
                'quiet': True,
                'no_warnings': True,
                'skip_download': True,
                'extract_flat': False,
                'no_check_certificate': True,
            }
            
            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                info = ydl.extract_info(video_url, download=False)
                
                title = info.get('title', 'Unknown Title')
                channel = info.get('uploader', 'Unknown Channel')
                duration = info.get('duration', 0)
                view_count = info.get('view_count', 0)
                upload_date = info.get('upload_date', '')
                
                # Format duration
                duration_str = format_duration(duration)
                
                # Format upload date
                formatted_date = ""
                if upload_date and len(upload_date) == 8:
                    try:
                        year = upload_date[:4]
                        month = upload_date[4:6]
                        day = upload_date[6:8]
                        formatted_date = f"{year}-{month}-{day}"
                    except:
                        formatted_date = upload_date
                
                return {
                    'title': title,
                    'channel': channel,
                    'duration': duration_str,
                    'thumbnail': f"https://img.youtube.com/vi/{video_id}/maxresdefault.jpg",
                    'transcript_text': transcript_text,
                    'video_id': video_id,
                    'view_count': view_count,
                    'upload_date': formatted_date,
                    'video_duration': source.video_duration if hasattr(source, 'video_duration') else duration_str,
                    'word_count': len(transcript_text.split()) if transcript_text else 0,
                    'source': 'database_transcript'
                }
                
        except Exception as ydl_error:
            print(f"yt-dlp error: {str(ydl_error)}")
            
            # Fallback to basic information
            return {
                'title': 'YouTube Video',
                'channel': 'Unknown Channel',
                'duration': source.video_duration if hasattr(source, 'video_duration') else 'Unknown',
                'thumbnail': f"https://img.youtube.com/vi/{video_id}/maxresdefault.jpg",
                'transcript_text': transcript_text,
                'video_id': video_id,
                'view_count': 0,
                'upload_date': '',
                'word_count': len(transcript_text.split()) if transcript_text else 0,
                'source': 'fallback',
                'note': 'Limited preview - video metadata unavailable'
            }
            
    except Exception as e:
        print(f"YouTube preview error: {str(e)}")
        raise Exception(f'Error getting YouTube preview: {str(e)}')

def extract_youtube_video_id(url):
    """Extract video ID from YouTube URL with improved pattern matching"""
    if not url:
        return None
        
    # Handle different YouTube URL formats
    patterns = [
        r'(?:youtube\.com\/watch\?v=|youtu\.be\/|youtube\.com\/embed\/|youtube\.com\/v\/)([a-zA-Z0-9_-]{11})',
        r'youtube\.com\/watch\?.*v=([a-zA-Z0-9_-]{11})',
    ]
    
    for pattern in patterns:
        match = re.search(pattern, url)
        if match:
            return match.group(1)
    
    return None

def format_duration(seconds):
    """Format duration from seconds to HH:MM:SS or MM:SS"""
    if not seconds or seconds == 0:
        return "Unknown"
    
    try:
        seconds = int(seconds)
        hours = seconds // 3600
        minutes = (seconds % 3600) // 60
        secs = seconds % 60
        
        if hours > 0:
            return f"{hours}:{minutes:02d}:{secs:02d}"
        else:
            return f"{minutes}:{secs:02d}"
    except (ValueError, TypeError):
        return "Unknown"